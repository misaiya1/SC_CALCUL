#!/usr/bin/env python
# _*_ coding:utf-8 _*_
from UI.noname import SC_CALCUL
import numpy as np
from math import pi
import datetime

import os
import sys
import matplotlib.pyplot as plt
from matplotlib.pyplot import Polygon
import wx
import wx.xrc
import math
from docx import Document
from docx.shared import Inches
from docx.shared import Pt
from docx.shared import RGBColor
from docx.oxml.ns import qn
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT

from docxtpl import DocxTemplate
from pptx import Presentation
from pptx_tools import utils
from pptx.util import Inches, Pt

myChecked = False
NetFreq = 50
Root2 = 1.414
Root3 = 1.732
TWOPI = 2 * pi
''' 函数：返回变量是否被定义过 '''


def isset(v):
    try:
        type(eval(v))
    except:
        return 0
    else:
        return 1


def func_cycle(x, e, x_o, y_o):
    return np.sqrt((e ** 2 - (x - x_o) ** 2)) + y_o


class MyFrame(SC_CALCUL):
    ''' 创建输出路径  '''
    global cwd
    cwd = os.getcwd()
    print(cwd)
    global outputFolderName

    outputFolderName = cwd + r'\fig_save'  # + datetime.datetime.now().strftime('%Y.%m.%d.%H.%M')
    outPathExists = os.path.exists(outputFolderName)
    if not outPathExists:
        os.makedirs(outputFolderName)
        print(outputFolderName + ' 路径创建成功')
    else:
        print(outputFolderName + ' 路径已存在')

    def __init__(self, parent):
        SC_CALCUL.__init__(self, parent)

        wx.Frame.SetBackgroundColour(self, 'Write')  # ???
        font1 = wx.Font(22, wx.MODERN, wx.NORMAL, wx.NORMAL)
        self.m_staticText15.SetFont(font1)
        self.m_button.SetFont(font1)
        self.m_button2.SetFont(font1)
        self.m_staticText16.SetFont(font1)

        self.m_grid2.SetRowLabelValue(0, '转速[RPM]')  # ???
        self.m_grid2.SetRowLabelValue(1, '载荷需求转矩（长期）[Nm]')  # ???
        self.m_grid2.SetRowLabelValue(2, '载荷需求转矩（短时）[Nm]')  # ???
        self.m_grid2.SetRowLabelValue(3, '磁链给定')  # ???
        # 第一列加宽
        # self.m_grid2.AutoSizeColLabelSize()
        # self.m_grid2.AutoSizeRowLabelSize()
        # self.m_grid2.AutoSize()
        # self.m_grid2.Fit()
        self.m_grid2.AutoSizeRowLabelSize(1)
        self.m_grid2.SetRowLabelSize(160)
        # ctrl - c

        ########################################################## 3.0 - 150
        self.m_grid2.SetCellValue(0, 0, '550')  # 第1行第0个
        self.m_grid2.SetCellValue(0, 1, '700')
        self.m_grid2.SetCellValue(0, 2, '800')
        self.m_grid2.SetCellValue(0, 3, '900')
        self.m_grid2.SetCellValue(0, 4, '1000')
        self.m_grid2.SetCellValue(0, 5, '1100')
        self.m_grid2.SetCellValue(0, 6, '1200')
        self.m_grid2.SetCellValue(0, 7, '1300')
        self.m_grid2.SetCellValue(0, 8, '1400')
        self.m_grid2.SetCellValue(0, 9, '1450')
        self.m_grid2.SetCellValue(0, 10, '1500')
        self.m_grid2.SetCellValue(0, 11, '1600')
        self.m_grid2.SetCellValue(0, 12, '1700')
        self.m_grid2.SetCellValue(0, 13, '1800')
        self.m_grid2.SetCellValue(0, 14, '1900')
        self.m_grid2.SetCellValue(0, 15, '2000')
        self.m_grid2.SetCellValue(0, 16, '2100')
        self.m_grid2.SetCellValue(0, 17, '2200')
        self.m_grid2.SetCellValue(0, 18, '2300')

        self.m_grid2.SetCellValue(1, 0, '-29700')  # 第1行第0个
        self.m_grid2.SetCellValue(1, 1, '-29700')
        self.m_grid2.SetCellValue(1, 2, '-29700')
        self.m_grid2.SetCellValue(1, 3, '-29700')
        self.m_grid2.SetCellValue(1, 4, '-29700')
        self.m_grid2.SetCellValue(1, 5, '-29700')
        self.m_grid2.SetCellValue(1, 6, '-29700')
        self.m_grid2.SetCellValue(1, 7, '-29700')
        self.m_grid2.SetCellValue(1, 8, '-29700')
        self.m_grid2.SetCellValue(1, 9, '-29700')
        self.m_grid2.SetCellValue(1, 10, '-29100')
        self.m_grid2.SetCellValue(1, 11, '-27300')
        self.m_grid2.SetCellValue(1, 12, '-25700')
        self.m_grid2.SetCellValue(1, 13, '-24250')
        self.m_grid2.SetCellValue(1, 14, '-22900')
        self.m_grid2.SetCellValue(1, 15, '-21700')
        self.m_grid2.SetCellValue(1, 16, '-20600')
        self.m_grid2.SetCellValue(1, 17, '-19550')
        self.m_grid2.SetCellValue(1, 18, '-18600')

        self.m_grid2.SetCellValue(2, 0, '-32950')  # 第1行第0个
        self.m_grid2.SetCellValue(2, 1, '-32950')
        self.m_grid2.SetCellValue(2, 2, '-32950')
        self.m_grid2.SetCellValue(2, 3, '-32950')
        self.m_grid2.SetCellValue(2, 4, '-32950')
        self.m_grid2.SetCellValue(2, 5, '-32950')
        self.m_grid2.SetCellValue(2, 6, '-32950')
        self.m_grid2.SetCellValue(2, 7, '-32950')
        self.m_grid2.SetCellValue(2, 8, '-32950')
        self.m_grid2.SetCellValue(2, 9, '-32950')
        self.m_grid2.SetCellValue(2, 10, '-32050')
        self.m_grid2.SetCellValue(2, 11, '-30050')
        self.m_grid2.SetCellValue(2, 12, '-28200')
        self.m_grid2.SetCellValue(2, 13, '-26550')
        self.m_grid2.SetCellValue(2, 14, '-25050')
        self.m_grid2.SetCellValue(2, 15, '-23700')
        self.m_grid2.SetCellValue(2, 16, '-22450')
        self.m_grid2.SetCellValue(2, 17, '-21300')
        self.m_grid2.SetCellValue(2, 18, '-20250')

        self.m_grid2.SetCellValue(3, 10, '1.90')
        self.m_grid2.SetCellValue(3, 11, '1.80')
        self.m_grid2.SetCellValue(3, 12, '1.70')
        self.m_grid2.SetCellValue(3, 13, '1.60')
        self.m_grid2.SetCellValue(3, 14, '1.50')
        self.m_grid2.SetCellValue(3, 15, '1.40')
        self.m_grid2.SetCellValue(3, 16, '1.35')
        self.m_grid2.SetCellValue(3, 17, '1.25')
        self.m_grid2.SetCellValue(3, 18, '1.20')

        self.statusbar = self.CreateStatusBar()
        self.Centre()

    def OnCheck(self, event):
        global myChecked
        w = 2 * pi * NetFreq

        myChecked = self.m_checkBox1.GetValue()
        if myChecked:
            self.m_Lksss.SetLabel(u"定子漏抗[ohm]")
            self.m_Lkrrr.SetLabel(u"转子漏抗[ohm]")
            self.m_staticText22.SetLabel(u"激磁电抗[ohm]")

            self.m_Lm.SetValue(str(float(self.m_Lm.GetValue()) * w))
            self.m_Lks.SetValue(str(float(self.m_Lks.GetValue()) * w))
            self.m_Lkr.SetValue(str(float(self.m_Lkr.GetValue()) * w))



        else:
            self.m_Lksss.SetLabel(u"定子漏感[H]")
            self.m_Lkrrr.SetLabel(u"转子漏感[H]")
            self.m_staticText22.SetLabel(u"电机互感[H]")

            self.m_Lm.SetValue(str(float(self.m_Lm.GetValue()) / w))
            self.m_Lks.SetValue(str(float(self.m_Lks.GetValue()) / w))
            self.m_Lkr.SetValue(str(float(self.m_Lkr.GetValue()) / w))

        print('myChecked   %f' % myChecked)

    def m_buttonOnButtonClick(self, event):
        print("press button1")
        self.DoWork(1)
        self.PlotAll()

    def m_buttonOnButtonClick2(self, event):
        print("press button2")  # 报告生成
        self.DoWork(1)
        self.Plot1()
        self.Plot2()
        self.Plot3()

        ###################################################################报错
        self.DocuReport()
        self.statusbar.SetStatusText('Report Generate at , %s ' % wx.Now())
        dlg = wx.MessageDialog(None, u"报告已生成，请确认", u"完成",
                               wx.YES_DEFAULT | wx.ICON_QUESTION)
        if dlg.ShowModal() == wx.ID_YES:
            self.Close(True)
        dlg.Destroy()
        ###################################################################不报错

        # try:
        #     self.DocuReport()
        #     self.statusbar.SetStatusText('Report Generate at , %s ' % wx.Now())
        #     dlg = wx.MessageDialog(None, u"报告已生成，请确认", u"完成", wx.YES_DEFAULT | wx.ICON_QUESTION)
        #     if dlg.ShowModal() == wx.ID_YES:
        #         self.Close(True)
        #     dlg.Destroy()
        # except:
        #     dlg = wx.MessageDialog(None, u"请先关闭生成物", u"完成", wx.YES_DEFAULT | wx.ICON_QUESTION)
        #     if dlg.ShowModal() == wx.ID_YES:
        #         self.Close(True)
        #     dlg.Destroy()

    def m_buttonOnButtonClick3(self, event):
        print("press button3")  # 无功边界
        self.DoWork(1)


    def mOnMenuSelection2(self, event):
        dlg = wx.MessageDialog(None,
                               "1.采用转子磁链定向，因此phirq=0，ird=0；\n2.电动机模式，主要涉及到数值计算，转矩可取正或负；\n3.鼠笼机的弱磁很简单，即减小定子侧的励磁电流，也就是重新调整此表中的V；\n4.然后再通过电流能力判断即可；\n5.网侧变流器按照单位功率因数计算;\n6.忽略摩擦损耗及铁耗，考虑铜耗！",
                               u"确认", wx.YES_DEFAULT | wx.ICON_QUESTION)
        if dlg.ShowModal() == wx.ID_YES:
            self.Close(True)
        dlg.Destroy()

    def mOnMenuSelection1(self, event):
        dlg = wx.MessageDialog(None, u"版本:V0.0", u"确认", wx.YES_DEFAULT | wx.ICON_QUESTION)
        if dlg.ShowModal() == wx.ID_YES:
            self.Close(True)
        dlg.Destroy()

    def m_slider1OnScrollChanged(self, event):
        pass
        # temp = float(self.m_slider1.GetValue())
        # self.m_textCtrlS.Value = ("%8.3f" % temp)

    def DoWork(self, netVotGain):

        plt.close('all')
        # slider input part
        # tempMin = float(self.m_textCtrlS1.GetValue())
        # tempMax = float(self.m_textCtrlS2.GetValue())
        # self.m_slider1.SetMin(tempMin)
        # self.m_slider1.SetMax(tempMax)
        # self.m_slider1.Value = temp
        # print(self.m_slider1.Value)

        # input part
        global pp, NetFreq, TongBuZhuanSu, StatorRePower, RotorOpenVoltage, GenRpm, XiaoLv, GridVRMS
        global Xm, Xs, Xr, Rs, Rr, Vpp, Xm__Xss, NetPF

        pp = float(self.m_pp.GetValue())  # pp
        A2 = pp

        GenRateFreq = float(self.m_edhz.GetValue())  # 电机额定频率
        B2 = GenRateFreq

        VdcRate = float(self.m_Vdc.GetValue())  # 额定直流母线电压
        C2 = VdcRate

        RateVolt = float(self.m_RateVolt.GetValue())  # 发电机额定电压
        D2 = RateVolt

        GridVRMS = float(self.m_GridVRMS.GetValue()) * netVotGain  # 电网电压
        D5 = GridVRMS

        NetPF = float(self.m_NewPF.GetValue())  # 电网电压

        RateGenCurrent = float(self.m_RateGenCurrent.GetValue())  # 发电机额定电流A
        E2 = RateGenCurrent

        XiaoLvGen = float(self.m_XiaoLvGen.GetValue())  # 发电机效率
        F2 = XiaoLvGen

        GenPower = float(self.m_GenPower.GetValue())  # 发电机额定功率W
        Vpp = GenPower
        ###########################################################################
        Rs = float(self.m_Rs.GetValue())  # 定子电阻
        H2 = Rs
        print('Rs = ')
        print(Rs)
        Rr = float(self.m_Rr.GetValue())  # 转子电阻
        I2 = Rr

        tempW = GenRateFreq * 2 * pi

        if not myChecked:
            Lm = float(self.m_Lm.GetValue())  # 电机互感
            L2 = Lm
            Lks = float(self.m_Lks.GetValue())  # 定子漏感
            J2 = Lks
            Lkr = float(self.m_Lkr.GetValue())  # 转子漏感
            K2 = Lkr
        else:
            Lm = float(self.m_Lm.GetValue()) / tempW  # 电机互感
            L2 = Lm
            Lks = float(self.m_Lks.GetValue()) / tempW  # 定子漏感
            J2 = Lks
            Lkr = float(self.m_Lkr.GetValue()) / tempW  # 转子漏感
            K2 = Lkr
        Xm = Lm * tempW
        Xs = Lks * tempW
        Xr = Lkr * tempW
        J5 = Xs
        K5 = Xr
        K5 = Xm

        Lss = Lm + Lks
        Lrr = Lm + Lkr
        Lamada = 1 - Lm ** 2 / Lss / Lrr
        J8 = Lss
        K8 = Lrr
        L8 = Lamada

        print('myChecked   HAVE %f' % myChecked)
        print('======================================')
        print('Lks  %f' % Lks)
        print('Lkr  %f' % Lkr)
        print('Lm  %f' % Lm)
        print('======================================')
        ############################################################################
        XiaoLv = float(self.m_XiaoLv.GetValue())  # 变流器效率
        M2 = XiaoLv

        NetFreq = float(self.m_hz.GetValue())  # 电网频率
        N2 = NetFreq

        SetGenCurRate = float(self.m_SetGenCurRate.GetValue())  # 机侧变流器额定电流 Arms
        M8 = SetGenCurRate

        SetGenCurMax = float(self.m_SetGenCurMax.GetValue())  # 机侧变流器极限电流 Arms
        M5 = SetGenCurMax

        SetNetCurRate = float(self.m_SetNetCurRate.GetValue())  # 网侧变流器额定电流 Arms
        M11 = SetNetCurRate

        SetNetCurMax = float(self.m_SetNetCurMax.GetValue())  # 网侧变流器极限电流 Arms
        M14 = SetNetCurMax

        SetSatMaxI = float(self.m_SetSatMaxI.GetValue())  # 定子极限电流 Arms
        N5 = SetSatMaxI

        # TongBuZhuanSu = 60 * NetFreq / pp  # 同步转速
        # B2 = TongBuZhuanSu
        # GoNetPower = float(self.m_GoNetPower.GetValue()) * 1E3  # 变流器容量
        # C2 = GoNetPower
        # StatorRePower = float(self.m_StatorRePower.GetValue()) * 1E3  # 定子无功
        # D2 = StatorRePower
        # RotorOpenVoltage = float(self.m_RotorOpen.GetValue()) * netVotGain # 转子开口电压
        # E2 = RotorOpenVoltage

        Vpp = Root2 * GridVRMS / Root3
        print(Vpp)
        GenRpm = float(self.m_textCtrlS.GetValue())  # 发电机转

        ###########################################################################

        global iGenRpm, iTorqueLongTime, iTorqueShortTime, iPsi, tempMin, tempMax, n
        tempMin = 1e6
        tempMax = 0
        tempLeft = 0
        tempLeft1 = 0
        tempLeft2 = 0
        tempLeft3 = 0

        tempRight = 1e6
        tempRight1 = 1e6
        tempRight2 = 1e6
        tempRight3 = 1e6

        tempLong = 0;
        tempShort = 0;
        tempPsi = 0;

        self.m_grid2.SetRowLabelValue(0, '转速[RPM]')  # ???
        self.m_grid2.SetRowLabelValue(1, '载荷需求转矩（长期）[Nm]')  # ???
        self.m_grid2.SetRowLabelValue(2, '载荷需求转矩（短时）[Nm]')  # ???
        self.m_grid2.SetRowLabelValue(3, '磁链给定')  # ???

        m = 26
        n = 0
        baohan = 0

        iGenRpm = []
        iTorqueLongTime = []
        iTorqueShortTime = []
        iPsi = []
        for i in range(m):
            if (self.m_grid2.GetCellValue(0, i) == ''):
                pass
                # print('empty %s' % i)
            else:
                # print('full %s' % i)
                temp1 = float(self.m_grid2.GetCellValue(0, i))
                temp2 = float(self.m_grid2.GetCellValue(1, i))
                temp3 = float(self.m_grid2.GetCellValue(2, i))
                try:
                    temp4 = float(self.m_grid2.GetCellValue(3, i))
                except:
                    temp4 = 0;

                iGenRpm.append(float(temp1))
                iTorqueLongTime.append(float(temp2) * 1E0)
                iTorqueShortTime.append(float(temp3) * 1E0)
                iPsi.append(float(temp4) * 1E0)

                tempMin = min(tempMin, float(self.m_grid2.GetCellValue(0, i)))
                tempMax = max(tempMax, float(self.m_grid2.GetCellValue(0, i)))

                if GenRpm == temp1:  # 包含
                    baohan = 1
                    tempLong = temp2
                    tempShort = temp3
                    tempPsi = temp4

                if temp1 < GenRpm and temp1 >= tempLeft:
                    tempLeft = temp1  # 刷新最小值
                    tempLeft1 = temp2
                    tempLeft2 = temp3
                    tempLeft3 = temp4

                if temp1 > GenRpm and temp1 <= tempRight:
                    tempRight = temp1  # 刷新最大值
                    tempRight1 = temp2
                    tempRight2 = temp3
                    tempRight3 = temp4

                n = n + 1
        print('num %s' % n)

        if GenRpm < tempMin:  # 比最小值还小？
            GenRpm = tempMin  # 不允许
            self.m_textCtrlS.SetValue(str(tempMin))

        if GenRpm > tempMax:
            GenRpm = tempMax
            print('tempMax %s' % tempMax)
            self.m_textCtrlS.SetValue(str(tempMax))

        # 队列对iGenRpm进行排序，用index排序iTorqueLongTime iTorqueShortTime
        iGenRpm = np.array(iGenRpm)
        iTorqueLongTime = np.array(iTorqueLongTime)
        iTorqueShortTime = np.array(iTorqueShortTime)
        iPsi = np.array(iPsi)

        index = np.argsort(iGenRpm)
        iGenRpm = iGenRpm[index]
        iTorqueLongTime = iTorqueLongTime[index]
        iTorqueShortTime = iTorqueShortTime[index]
        iPsi = iPsi[index]

        # 队列最后加入 单点值
        # iGenRpm.append(GenRpm)
        ###################################
        iGenRpm = np.append(iGenRpm, GenRpm)
        if baohan == 1:
            # iTorqueLongTime.append(tempLong * 1E3)
            iTorqueLongTime = np.append(iTorqueLongTime, tempLong * 1E0)
            iTorqueShortTime = np.append(iTorqueShortTime, tempShort * 1E0)
            iPsi = np.append(iPsi, tempPsi * 1E0)
        else:
            # 取左右2个转速，按比例增加iTorqueLongTime iTorqueShortTime
            tempLong = (GenRpm - tempLeft) / (tempRight - tempLeft) * (tempRight1 - tempLeft1) + tempLeft1
            iTorqueLongTime = np.append(iTorqueLongTime, tempLong * 1E0)

            tempShort = (GenRpm - tempLeft) / (tempRight - tempLeft) * (tempRight2 - tempLeft2) + tempLeft2
            iTorqueShortTime = np.append(iTorqueShortTime, tempShort * 1E0)

            tempPsi = (GenRpm - tempLeft) / (tempRight - tempLeft) * (tempRight3 - tempLeft3) + tempLeft3
            iPsi = np.append(iPsi, tempPsi * 1E0)

        n = n + 1
        ##########################################################################作图
        global iO2, iP2, iQ2, iR2, iS2, iT2, iU2, iV2, iW2, iX2, iY2, iZ2
        global iAX2, iAY2, iAZ2, iBA2, iBB2, iBC2, iBD2

        iO2 = [0 for i in range(n)]
        iO2 = iGenRpm
        iP2 = [0 for i in range(n)]  # 扭矩长期
        iP2 = iTorqueLongTime
        iQ2 = [0 for i in range(n)]  # 发电机轴功率
        iR2 = [0 for i in range(n)]
        iS2 = [0 for i in range(n)]  # 扭矩短期
        iS2 = iTorqueShortTime
        iT2 = [0 for i in range(n)]
        iU2 = [0 for i in range(n)]
        iV2 = [0 for i in range(n)]
        iW2 = [0 for i in range(n)]
        iX2 = [0 for i in range(n)]
        iY2 = [0 for i in range(n)]
        iZ2 = [0 for i in range(n)]
        iAA2 = [0 for i in range(n)]
        iAB2 = [0 for i in range(n)]
        iAC2 = [0 for i in range(n)]
        iAD2 = [0 for i in range(n)]
        iAE2 = [0 for i in range(n)]
        iAF2 = [0 for i in range(n)]
        iAG2 = [0 for i in range(n)]
        iAH2 = [0 for i in range(n)]
        iAI2 = [0 for i in range(n)]
        iAJ2 = [0 for i in range(n)]
        iAK2 = [0 for i in range(n)]
        iAL2 = [0 for i in range(n)]
        iAM2 = [0 for i in range(n)]
        iAN2 = [0 for i in range(n)]
        iAO2 = [0 for i in range(n)]
        iAP2 = [0 for i in range(n)]
        iAQ2 = [0 for i in range(n)]
        iAR2 = [0 for i in range(n)]
        iAS2 = [0 for i in range(n)]
        iAT2 = [0 for i in range(n)]
        iAU2 = [0 for i in range(n)]
        iAV2 = [0 for i in range(n)]
        iAW2 = [0 for i in range(n)]
        iAX2 = [0 for i in range(n)]
        iAY2 = [0 for i in range(n)]
        iAZ2 = [0 for i in range(n)]
        iBA2 = [0 for i in range(n)]
        iBB2 = [0 for i in range(n)]
        iBC2 = [0 for i in range(n)]
        iBD2 = [0 for i in range(n)]

        # O:RPM-- P:Long-- Q:Pmlong-- R:PeLongAD--- S:Short--  T:Pmshort--  U:PeShortAE---  V:psi--  W:isd0--  X:isqLong--  Y:isqShort--  Z:irqLong--
        # AA:isqShort--  AB:zcjplLong--  AC:zcjplShort--  AD:WgenLong--  AE:WgenShort--  AF:zclLong--  AG:zclShort--  AH:GenHzLong--  AI:GenHzShort--
        # AJ:usdLong--  AK:usqLong--  AL:usLong--  AM:pj1  AN:usdShort--  AO:usqShort--  AP:usShort--  AQ:pj2  AR:--
        # AS:  AT:  AU:  AV:  AW:  AX:  AY:  AZ:

        for i in range(n):
            iQ2[i] = iO2[i] * 2 * 3.14 / 60 * iP2[i]
            iT2[i] = iO2[i] * 2 * 3.14 / 60 * iS2[i]

            iV2[i] = D2 / 1.732 * 1.414 / 2 / pi / B2
            if iPsi[i] != 0:
                iV2[i] = iPsi[i]

            iW2[i] = iV2[i] / Lm
            iX2[i] = 2 * iP2[i] * Lrr / 3 / pp / Lm / iV2[i]
            iY2[i] = 2 * iS2[i] * Lrr / 3 / pp / Lm / iV2[i]
            iZ2[i] = -Lm / Lrr * iX2[i]

            iAA2[i] = -Lm / Lrr * iY2[i]
            iAB2[i] = Rr / Lrr * iX2[i] / iW2[i]
            iAC2[i] = Rr / Lrr * iY2[i] / iW2[i]
            iAD2[i] = iAB2[i] + iO2[i] * 2 * 3.14 / 60 * pp
            iAE2[i] = iAC2[i] + iO2[i] * 2 * 3.14 / 60 * pp

            iR2[i] = iP2[i] * iAD2[i] / pp
            iU2[i] = iS2[i] * iAE2[i] / pp

            iAF2[i] = (iAD2[i] - 6.28 / 60 * iO2[i] * pp) / iAD2[i]
            iAG2[i] = (iAE2[i] - 6.28 / 60 * iO2[i] * pp) / iAE2[i]
            iAH2[i] = iAD2[i] / 2 / 3.14
            iAI2[i] = iAE2[i] / 2 / 3.14

            iAJ2[i] = Rs * iW2[i] - iAD2[i] * Lss * Lamada * iX2[i]
            iAK2[i] = Rs * iX2[i] + iAD2[i] * Lss * iW2[i]
            iAL2[i] = (iAJ2[i] ** 2 + iAK2[i] ** 2) ** 0.5 / 1.414

            iAN2[i] = Rs * iW2[i] - iAE2[i] * (Lss * iY2[i] + Lm * iAA2[i])  # 歧义
            iAO2[i] = Rs * iY2[i] + iAE2[i] * Lss * iW2[i]
            iAP2[i] = (iAN2[i] ** 2 + iAO2[i] ** 2) ** 0.5 / 1.414

            iAM2[i] = 1 if (iAL2[i] < VdcRate / 1.732 / 1.414) else 0
            iAQ2[i] = 1 if (iAP2[i] < VdcRate / 1.732 / 1.414) else 0

            iAS2[i] = 1.5 * (iAJ2[i] * iW2[i] + iX2[i] * iAK2[i])
            iAT2[i] = iAS2[i] - iR2[i]
            iAU2[i] = 3 * ((iW2[i] ** 2 + iX2[i] ** 2) ** 0.5 / 1.414) ** 2 * Rs
            iAV2[i] = iAF2[i] * iR2[i]
            iAW2[i] = 3 * (iZ2[i] / 1.414) ** 2 * Rr
            iAX2[i] = 1 - (iAU2[i] + iAW2[i]) / abs(iQ2[i])
            iAY2[i] = abs(iAS2[i]) / ((iW2[i] ** 2 + iX2[i] ** 2) ** 0.5 / 1.414 * iAL2[i] * 3)

            iAZ2[i] = (iW2[i] ** 2 + iX2[i] ** 2) ** 0.5 / 1.414
            iBA2[i] = (iW2[i] ** 2 + iY2[i] ** 2) ** 0.5 / 1.414

            CVTeffi = XiaoLv
            Ug = GridVRMS

            iBB2[i] = abs((iR2[i] * CVTeffi) / Ug / 1.732)
            iBC2[i] = abs(iBB2[i] / NetPF)
            iBD2[i] = (1 - NetPF ** 2) ** 0.5 / NetPF * iR2[i]

            iQ2[i] = iQ2[i] * 1E-3
            iR2[i] = iR2[i] * 1E-3
            iT2[i] = iT2[i] * 1E-3
            iU2[i] = iU2[i] * 1E-3

            iBD2[i] = iBD2[i] * 1E-3

        self.m_V.Value = ("%.3f" % iV2[-1])
        self.m_W.Value = ("%.3f" % iW2[-1])
        self.m_Q.Value = ("%.3f" % (iQ2[-1] * 1E0))
        self.m_R.Value = ("%.3f" % (iR2[-1] * 1E0))
        self.m_X.Value = ("%.3f" % iX2[-1])
        self.m_Z.Value = ("%.3f" % iZ2[-1])
        self.m_AB.Value = ("%.3f" % iAB2[-1])
        self.m_AD.Value = ("%.3f" % iAD2[-1])
        self.m_AF.Value = ("%.5f" % iAF2[-1])
        self.m_AH.Value = ("%.3f" % iAH2[-1])
        self.m_AJ.Value = ("%.3f" % iAJ2[-1])
        self.m_AK.Value = ("%.3f" % iAK2[-1])
        self.m_AL.Value = ("%.3f" % iAL2[-1])
        self.m_AM.Value = ("%.3f" % iAM2[-1])

        self.m_T.Value = ("%.3f" % (iT2[-1] * 1E0))
        self.m_U.Value = ("%.3f" % (iU2[-1] * 1E0))
        self.m_Y.Value = ("%.3f" % iY2[-1])
        self.m_AA.Value = ("%.3f" % iAA2[-1])
        self.m_AC.Value = ("%.3f" % iAC2[-1])
        self.m_AE.Value = ("%.3f" % iAE2[-1])
        self.m_AG.Value = ("%.5f" % iAG2[-1])
        self.m_AI.Value = ("%.3f" % iAI2[-1])
        self.m_AN.Value = ("%.3f" % iAN2[-1])
        self.m_AO.Value = ("%.3f" % iAO2[-1])
        self.m_AP.Value = ("%.3f" % iAP2[-1])
        self.m_AQ.Value = ("%.3f" % iAQ2[-1])

        self.m_AS.Value = ("%.3f" % (iAS2[-1] * 1E0))
        self.m_AT.Value = ("%.3f" % (iAT2[-1] * 1E0))
        self.m_AU.Value = ("%.3f" % (iAU2[-1] * 1E0))
        self.m_AV.Value = ("%.3f" % (iAV2[-1] * 1E0))
        self.m_AW.Value = ("%.3f" % (iAW2[-1] * 1E0))
        self.m_AX.Value = ("%.4f" % iAX2[-1])
        self.m_AY.Value = ("%.4f" % iAY2[-1])
        self.m_AZ.Value = ("%.4f" % iAZ2[-1])

        self.m_BB.Value = ("%.4f" % iBB2[-1])
        self.m_BC.Value = ("%.4f" % iBC2[-1])
        self.m_BD.Value = ("%.4f" % iBD2[-1])

        print("------------")

        ######################################################################################

    def PlotAll(self):

        plt.figure(5)
        gongLvTu = plt.subplot(2, 2, 1)  # 功率图
        gongLvTu.grid(linestyle='-.', which='major')
        gongLvTu.set_title('Scope - Power')

        miloc = plt.MultipleLocator(50)
        maloc = plt.MultipleLocator(100)
        gongLvTu.xaxis.set_minor_locator(miloc)
        gongLvTu.yaxis.set_minor_locator(maloc)
        gongLvTu.grid(linestyle='-.', which='minor', linewidth=0.3, alpha=0.9)

        iGenRpm_1 = iGenRpm[0:-1]
        iQ2_1 = iQ2[0:-1]  # 轴功率长期
        iR2_1 = iR2[0:-1]  # 电磁功率长期
        iT2_1 = iT2[0:-1]
        iU2_1 = iU2[0:-1]

        plt.plot(iGenRpm_1, iQ2_1, color="red", linewidth=1.5, linestyle="-", label="Shaft Power Long Time")
        plt.plot(iGenRpm_1, iR2_1, color="gray", linewidth=1.5, linestyle=":", label="Elec. Power Long Time")
        plt.plot(iGenRpm_1, iT2_1, color="blue", linewidth=1.5, linestyle="-", label="Shaft Power Short Time")
        plt.plot(iGenRpm_1, iU2_1, color="gray", linewidth=1.5, linestyle=":", label="Elec. Power Short Time")

        plt.xlim(tempMin - 50, tempMax + 50)
        # gongLvTu.set_xlabel('generetor speed[rpm]')
        gongLvTu.set_ylabel('[kW]')
        gongLvTu.legend()
        print(iGenRpm_1)
        print("tempMin = %s" % tempMin)
        print("tempMax = %s" % tempMax)
        # 标注额定点
        for i in range(n - 1):
            gongLvTu.scatter(iGenRpm_1[i], (iQ2_1[i]), s=30, color='gray')
            gongLvTu.scatter(iGenRpm_1[i], (iR2_1[i]), s=30, color='gray')
            gongLvTu.scatter(iGenRpm_1[i], (iT2_1[i]), s=30, color='gray')
            gongLvTu.scatter(iGenRpm_1[i], (iU2_1[i]), s=30, color='gray')

        # ######################################################################################
        dianLiuTu = plt.subplot(2, 2, 2)  # 电流图
        dianLiuTu.set_title('Scope - Current')
        dianLiuTu.grid(linestyle='-.', which='major')

        miloc = plt.MultipleLocator(50)
        maloc = plt.MultipleLocator(100)
        dianLiuTu.xaxis.set_minor_locator(miloc)
        dianLiuTu.yaxis.set_minor_locator(maloc)
        dianLiuTu.grid(linestyle='-.', which='minor', linewidth=0.3, alpha=0.9)

        iAZ_1 = iAZ2[0:-1]
        iBA2_1 = iBA2[0:-1]
        iBC2_1 = iBC2[0:-1]

        plt.plot(iGenRpm_1, iAZ_1, color="red", linewidth=1.5, linestyle="-", label="Istator RMS Long Time")  # 定子电流
        plt.plot(iGenRpm_1, iBA2_1, color="blue", linewidth=1.5, linestyle="--", label="Istator RMS Short TIme")
        plt.plot(iGenRpm_1, iBC2_1, color="pink", linewidth=1.5, linestyle="-.", label="Igrid RMS")

        plt.xlim(tempMin - 50, tempMax + 50)
        # dianLiuTu.set_xlabel('generetor speed[rpm]')
        dianLiuTu.set_ylabel('[A]')
        plt.legend()
        # 标注额定点
        for i in range(n - 1):
            dianLiuTu.scatter(iGenRpm_1[i], (iAZ_1[i]), s=30, color='gray')
            dianLiuTu.scatter(iGenRpm_1[i], (iBA2[i]), s=30, color='gray')
            dianLiuTu.scatter(iGenRpm_1[i], (iBC2[i]), s=30, color='gray')

        # ######################################################################################
        zhuanJuTu = plt.subplot(2, 2, 3)  #
        zhuanJuTu.set_title('Scope - Torque')
        zhuanJuTu.grid(linestyle='-.', which='major')

        miloc = plt.MultipleLocator(50)
        maloc = plt.MultipleLocator(1000)
        zhuanJuTu.xaxis.set_minor_locator(miloc)
        zhuanJuTu.yaxis.set_minor_locator(maloc)
        zhuanJuTu.grid(linestyle='-.', which='minor', linewidth=0.3, alpha=0.9)

        iP2_1 = iP2[0:-1]
        iS2_1 = iS2[0:-1]

        plt.plot(iGenRpm_1, iP2_1, color="red", linewidth=1.5, linestyle="-", label="Torque Require Long Time")  # 长期
        plt.plot(iGenRpm_1, iS2_1, color="blue", linewidth=1.5, linestyle="-.", label="Torque Require Short Time")  # 期

        plt.xlim(tempMin - 50, tempMax + 50)
        zhuanJuTu.set_xlabel('generetor speed[rpm]')
        zhuanJuTu.set_ylabel('[N.m]')
        plt.legend()
        # 标注额定点
        for i in range(n - 1):
            zhuanJuTu.scatter(iGenRpm_1[i], (iP2_1[i]), s=30, color='gray')
            zhuanJuTu.scatter(iGenRpm_1[i], (iS2_1[i]), s=30, color='gray')

        # ######################################################################################
        dianYaTu = plt.subplot(2, 2, 4)  # 效率图
        dianYaTu.set_title('Scope - Efficiency & Power Factor')
        dianYaTu.grid(linestyle='-.', which='major')

        miloc = plt.MultipleLocator(50)
        maloc = plt.MultipleLocator(100)
        dianYaTu.xaxis.set_minor_locator(miloc)
        dianYaTu.yaxis.set_minor_locator(maloc)
        dianYaTu.grid(linestyle='-.', which='minor', linewidth=0.3, alpha=0.9)

        iAX2_1 = iAX2[0:-1]
        iAY2_1 = iAY2[0:-1]

        plt.plot(iGenRpm_1, iAX2_1, color="red", linewidth=1.5, linestyle="-", label="Gen Effic.")  # 电机效率
        plt.plot(iGenRpm_1, iAY2_1, color="blue", linewidth=1.5, linestyle="-", label="Power Factor")  # 电机效率

        plt.xlim(tempMin - 50, tempMax + 50)
        dianYaTu.set_xlabel('generetor speed[rpm]')
        dianYaTu.set_ylabel('[%]')
        plt.legend()
        # 标注额定点
        for i in range(n - 1):
            dianYaTu.scatter(iGenRpm[i], (iAX2_1[i]), s=30, color='gray')
            dianYaTu.scatter(iGenRpm_1[i], (iAY2_1[i]), s=30, color='gray')

        plt.tight_layout()
        plt.savefig("fig1_4" + ".png", dpi=1000)
        plt.show()
        plt.close('all')

    def Plot1(self):
        pass
        # iGenRpm_1 = iGenRpm[0:-1]
        # plt.figure(1)
        # gongLvTu = plt.subplot(1, 1, 1)  # 功率图
        # gongLvTu.grid(linestyle='-.', which='major')
        # gongLvTu.set_title('Scope - Power')
        #
        # miloc = plt.MultipleLocator(50)
        # maloc = plt.MultipleLocator(100)
        # gongLvTu.xaxis.set_minor_locator(miloc)
        # gongLvTu.yaxis.set_minor_locator(maloc)
        # gongLvTu.grid(linestyle='-.', which='minor', linewidth=0.3, alpha=0.9)
        #
        # iStatorAPower_1 = iStatorAPower[0:-1]
        # iRotorAPower_1 = iRotorAPower[0:-1]
        # iTotalAPower_1 = iTotalAPower[0:-1]
        #
        # plt.plot(iGenRpm_1, iStatorAPower_1, color="red", linewidth=1.5, linestyle="-", label="Stator Active Power")
        # plt.plot(iGenRpm_1, iRotorAPower_1, color="blue", linewidth=1.5, linestyle="-.", label="Rotor Active Power")
        # plt.plot(iGenRpm_1, iTotalAPower_1, color="gray", linewidth=1.5, linestyle="--", label="Total Active Power")
        #
        # plt.xlim(tempMin - 50, tempMax + 50)
        # gongLvTu.set_xlabel('generetor speed[rpm]')
        # gongLvTu.set_ylabel('[kW]')
        # gongLvTu.legend()
        #
        # # 标注额定点
        # for i in range(n - 1):
        #     gongLvTu.scatter(iGenRpm_1[i], (iStatorAPower_1[i]), s=30, color='gray')
        #     gongLvTu.scatter(iGenRpm_1[i], (iRotorAPower_1[i]), s=30, color='gray')
        #     gongLvTu.scatter(iGenRpm_1[i], (iTotalAPower_1[i]), s=30, color='gray')
        #
        # plt.tight_layout()
        #
        # plt.savefig(outputFolderName + r"\fig_power.png", dpi=1000)
        # plt.close(1)

    def Plot2(self):
        pass
        # iGenRpm_1 = iGenRpm[0:-1]
        # plt.figure(2)
        # dianYaTu = plt.subplot(1, 1, 1)  # 电压图
        # dianYaTu.set_title('Scope - Rotor Voltage')
        # dianYaTu.grid(linestyle='-.', which='major')
        #
        # miloc = plt.MultipleLocator(50)
        # maloc = plt.MultipleLocator(100)
        # dianYaTu.xaxis.set_minor_locator(miloc)
        # dianYaTu.yaxis.set_minor_locator(maloc)
        # dianYaTu.grid(linestyle='-.', which='minor', linewidth=0.3, alpha=0.9)
        #
        # iRotorVoltRMS_1 = iRotorVoltRMS[0:-1]
        # iGenVrms_1 = iGenVrms[0:-1]
        #
        # # plt.plot(iGenRpm, iRotorVoltRMS, color="red", linewidth=1.5, linestyle="-", label="Vr RMS (estimate)")  # 定子电流
        # plt.plot(iGenRpm_1, iGenVrms_1, color="red", linewidth=1.5, linestyle="-", label="Vinv RMS")  # 定子电流
        #
        # plt.xlim(tempMin - 50, tempMax + 50)
        # dianYaTu.set_xlabel('generetor speed[rpm]')
        # dianYaTu.set_ylabel('[V]')
        # plt.legend()
        # # 标注额定点
        # for i in range(n - 1):
        #     # dianYaTu.scatter(iGenRpm[i], (iRotorVoltRMS[i]), s=30, color='gray')
        #     dianYaTu.scatter(iGenRpm_1[i], (iGenVrms_1[i]), s=30, color='gray')
        #
        # plt.tight_layout()
        # plt.savefig(outputFolderName + r"\fig_voltageR.png", dpi=1000)
        # plt.close(2)

    def Plot3(self):
        pass
        # iGenRpm_1 = iGenRpm[0:-1]
        # plt.figure(3)
        # dianLiuTu = plt.subplot(1, 1, 1)  # 电流
        # dianLiuTu.set_title('Scope - Current')
        # dianLiuTu.grid(linestyle='-.', which='major')
        #
        # miloc = plt.MultipleLocator(50)
        # maloc = plt.MultipleLocator(100)
        # dianLiuTu.xaxis.set_minor_locator(miloc)
        # dianLiuTu.yaxis.set_minor_locator(maloc)
        # dianLiuTu.grid(linestyle='-.', which='minor', linewidth=0.3, alpha=0.9)
        #
        # iStatorIrms_1 = iStatorIrms[0:-1]
        # iGenIrmsSJ_1 = iGenIrmsSJ[0:-1]
        # iNetIrms_1 = iNetIrms[0:-1]
        # iGridrms_1 = iGridrms[0:-1]
        #
        # plt.plot(iGenRpm_1, iStatorIrms_1, color="red", linewidth=1.5, linestyle="-", label="Istator RMS")  # 定子电流
        # plt.plot(iGenRpm_1, iGenIrmsSJ_1, color="blue", linewidth=1.5, linestyle="--", label="Irotor RMS")  # 机侧电流
        # plt.plot(iGenRpm_1, iNetIrms_1, color="gray", linewidth=1.5, linestyle="-.", label="Ipfc RMS")  # 网侧电流
        # plt.plot(iGenRpm_1, iGridrms_1, color="green", linewidth=1.5, linestyle=":", label="Igrid_total RMS")  # 上网电流
        #
        # plt.xlim(tempMin - 50, tempMax + 50)
        # dianLiuTu.set_xlabel('generetor speed[rpm]')
        # dianLiuTu.set_ylabel('[A]')
        # plt.legend()
        # # 标注额定点
        # for i in range(n - 1):
        #     dianLiuTu.scatter(iGenRpm_1[i], (iStatorIrms_1[i]), s=30, color='gray')
        #     dianLiuTu.scatter(iGenRpm_1[i], (iGenIrmsSJ_1[i]), s=30, color='gray')
        #     dianLiuTu.scatter(iGenRpm_1[i], (iNetIrms_1[i]), s=30, color='gray')
        #     dianLiuTu.scatter(iGenRpm_1[i], (iGridrms_1[i]), s=30, color='gray')
        #
        # plt.tight_layout()
        # plt.savefig(outputFolderName + r"\fig_Current.png", dpi=1000)
        # plt.close(3)



    def DocuReport(self):

        # '''save %0 data'''
        # iTotalAPower100 = iTotalAPower
        # iStatorAPower100 = iStatorAPower
        # iRotorAPower100 = iRotorAPower
        # iStatorIrms100 = list(map(abs, iStatorIrms))
        # iGenIrmsSJ100 = list(map(abs, iGenIrmsSJ))
        # iNetIrms100 = list(map(abs, iNetIrms))
        # iGridrms100 = list(map(abs, iGridrms))
        # iGenVrms100 = iGenVrms
        #
        # self.DoWork(1.1)
        # iTotalAPower110 = iTotalAPower
        # iStatorAPower110 = iStatorAPower
        # iRotorAPower110 = iRotorAPower
        # iStatorIrms110 = list(map(abs, iStatorIrms))
        # iGenIrmsSJ110 = list(map(abs, iGenIrmsSJ))
        # iNetIrms110 = list(map(abs, iNetIrms))
        # iGridrms110 = list(map(abs, iGridrms))
        # iGenVrms110 = iGenVrms
        #
        # self.DoWork(0.9)
        # iTotalAPower90 = iTotalAPower
        # iStatorAPower90 = iStatorAPower
        # iRotorAPower90 = iRotorAPower
        # iStatorIrms90 = list(map(abs, iStatorIrms))
        # iGenIrmsSJ90 = list(map(abs, iGenIrmsSJ))
        # iNetIrms90 = list(map(abs, iNetIrms))
        # iGridrms90 = list(map(abs, iGridrms))
        # iGenVrms90 = iGenVrms
        #
        self.DoWork(1)
        #####################################################################创建PPT,修改文本框
        print('ppt start')
        pptx = Presentation('abcd.pptx')
        for slide in pptx.slides:
            # 遍历幻灯片页的所有形状
            print("slide found")
            for shape in slide.shapes:
                print("shape found")
                # 判断形状是否含有文本框，如果含有则顺序运行代码
                if shape.has_text_frame:

                    # 获取文本框
                    text_frame = shape.text_frame
                    # 遍历文本框中的所有段落
                    for paragraph in text_frame.paragraphs:
                        paragraph.text = paragraph.text.replace('X1', str(format(Xs, '.5f')) + '[ohm]')
                        paragraph.text = paragraph.text.replace('X2', str(format(Xr, '.5f')) + '[ohm]')
                        paragraph.text = paragraph.text.replace('R1', str(format(Rs, '.5f')) + '[ohm]')
                        paragraph.text = paragraph.text.replace('R2', str(format(Rr, '.5f')) + '[ohm]')

                        paragraph.text = paragraph.text.replace('XNN', str(format(Xm, '.5f')) + '[ohm]')
                        paragraph.font.size = Pt(22)

        #####################################################################保存PPT,另存为png
        save_ppt = outputFolderName + r'\fig_x.pptx'
        pptx.save(save_ppt)

        utils.save_pptx_as_png(outputFolderName, save_ppt, overwrite_folder=True)
        fig_x_dir = outputFolderName + r'\幻灯片1.PNG'
        print('ppt end')
        #####################################################################创建word
        document = Document()
        # print(dir(document))

        document.styles['Normal'].font.name = '楷体'

        document.styles['Normal']._element.rPr.rFonts.set(qn('w:eastAsia'), u'楷体')
        # document.styles['Normal'].font.size = Pt(12)
        document.add_heading('鼠笼风电机组电驱系统设计计算报告', 0)

        p = document.add_paragraph('生成时间   ' + datetime.datetime.now().strftime('%Y.%m.%d  %H:%M'))

        document.add_heading('输入-电机参数', level=1)

        ###################参数表######################
        table_para = document.add_table(rows=17, cols=2)
        # print(dir(table_para))
        hdr_cells = table_para.columns[0].cells

        hdr_cells[0].text = ''
        hdr_cells[1].text = '极对数'
        hdr_cells[2].text = '发电机额定功率[kW]'
        hdr_cells[3].text = '电网频率[Hz]'
        hdr_cells[4].text = '额定频率[Hz]'
        hdr_cells[5].text = '变流器效率'
        hdr_cells[6].text = '电机效率'
        hdr_cells[7].text = '发电机额定电压[V]'
        hdr_cells[8].text = '电网电压[V]'
        hdr_cells[9].text = '上网功率因数'
        hdr_cells[10].text = '(选填)发电机额定电流[A]'
        hdr_cells[11].text = '(选填)额定直流母线电压[V]'
        hdr_cells[12].text = '(选填)网侧变流器额定电流[A]'
        hdr_cells[13].text = '(选填)网侧变流器极限电流[A]'
        hdr_cells[14].text = '(选填)机侧变流器额定电流[A]'
        hdr_cells[15].text = '(选填)机侧变流器极限电流[A]'
        hdr_cells[16].text = '(选填)发电机定子极限电流[A]'


        para_cells = table_para.columns[1].cells
        para_cells[0].text = ''
        para_cells[1].text = self.m_pp.GetValue()
        para_cells[2].text = self.m_GenPower.GetValue()
        para_cells[3].text = self.m_hz.GetValue()
        para_cells[4].text = self.m_edhz.GetValue()
        para_cells[5].text = self.m_XiaoLv.GetValue()
        para_cells[6].text = self.m_XiaoLvGen.GetValue()
        para_cells[7].text = self.m_RateVolt.GetValue()
        para_cells[8].text = self.m_GridVRMS.GetValue()
        para_cells[9].text = self.m_NewPF.GetValue()
        para_cells[10].text = self.m_RateGenCurrent.GetValue()
        para_cells[11].text = self.m_Vdc.GetValue()
        para_cells[12].text = self.m_SetNetCurRate.GetValue()
        para_cells[13].text = self.m_SetNetCurMax.GetValue()
        para_cells[14].text = self.m_SetGenCurRate.GetValue()
        para_cells[15].text = self.m_SetGenCurMax.GetValue()
        para_cells[16].text = self.m_SetSatMaxI.GetValue()
        #
        # ###################电机表######################
        # table_para = document.add_table(rows=5, cols=2)
        # # print(dir(table_para))
        # hdr_cells = table_para.columns[0].cells
        #
        # # self.m_Lksss.SetLabel(u"定子漏抗[ohm]")
        # # self.m_Lkrrr.SetLabel(u"转子漏抗[ohm]")
        # # self.m_staticText22.SetLabel(u"激磁电抗[ohm]")
        #
        # hdr_cells[0].text = self.m_staticText22.GetLabel()
        # hdr_cells[1].text = '定子电组[ohm]'
        # hdr_cells[2].text = '转子电阻[ohm]'
        # hdr_cells[3].text = self.m_Lksss.GetLabel()
        # hdr_cells[4].text = self.m_Lkrrr.GetLabel()
        #
        # para_cells = table_para.columns[1].cells
        # para_cells[0].text = self.m_Lm.GetValue()
        # para_cells[1].text = self.m_Rs.GetValue()
        # para_cells[2].text = self.m_Rr.GetValue()
        # para_cells[3].text = self.m_Lks.GetValue()
        # para_cells[4].text = self.m_Lkr.GetValue()
        #
        # ########等效电路图
        # try:
        #     paragraph = document.add_paragraph()
        #     paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        #     run = paragraph.add_run("")
        #     run.add_picture(r'.\fig_save\幻灯片1.PNG', width=Inches(5.25))
        #
        #     run = document.add_paragraph('图1. 鼠笼/发电机等效电路图')
        #     run.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        #
        # except:
        #     pass
        #
        # ########功率表
        # document.add_heading('输入-功率曲线', level=1)
        # n = len(iGenRpm) - 1
        # table = document.add_table(rows=1, cols=3)
        # hdr_cells = table.rows[0].cells
        # hdr_cells[0].text = '转速[rpm]'
        # hdr_cells[1].text = '上网功率[kW]'
        # hdr_cells[2].text = ' '
        #
        # for i in range(n):
        #     row_cells = table.add_row().cells
        #     row_cells[0].text = str(iGenRpm[i])
        #
        #     row_cells[1].text = str(format(iTorqueLongTime[i] * 0.001, '.2f'))
        #     row_cells[2].text = ' '
        #
        # ########输出功率
        # document.add_heading('输出-功率计算', level=1)
        #
        # paragraph = document.add_paragraph()
        # paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        # run = paragraph.add_run("")
        # run.add_picture(r'.\fig_save\fig_power.png', width=Inches(5.25))
        #
        # run = document.add_paragraph('图2. 定子/转子/总上网功率')
        # run.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        # run = document.add_paragraph('')
        # run.add_run(r'"*注：对于转子功率(Rotor Active Power)，负值表示从电网获取功率"').italic = True
        # document.add_paragraph('总上网功率最大值为：\t%8.3f [kW]' % (max(iTotalAPower100)))
        # document.add_paragraph('定子功率最大值为：  \t%8.3f [kW]' % (max(iStatorAPower100)))
        # document.add_paragraph('转子功率最大值为：  \t%8.3f [kW]' % (max(iRotorAPower100)))
        # document.add_paragraph('转子功率最小值为：  \t%8.3f [kW]' % (min(iRotorAPower100)))
        # document.add_page_break()
        #
        # ########输出电流
        # document.add_heading('输出-电流', level=1)
        #
        # paragraph = document.add_paragraph()
        # paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        # run = paragraph.add_run("")
        # run.add_picture(r'.\fig_save\fig_Current.png', width=Inches(5.25))
        #
        # run = document.add_paragraph('图3. 定子/转子/变流器网侧/总上网电流')
        # run.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        # run = document.add_paragraph('')
        # run.add_run(r'"*注：对于变流器网侧电流(Ipfc)，负值表示流向电网"').italic = True
        # run.add_run('\n').italic = True
        # run.add_run(r'"*注：转子电流、变流器网侧电流均为实际值，而非归算值"').italic = True
        #
        # document.add_paragraph(r'以下三个数据分别表示工况（正常电网电压）（电网90%低电压）（电网110%高电压）')
        # document.add_paragraph('定子电流最大值为：      \t%8.3f [A],   %8.3f [A],   %8.3f [A]' % (
        #     max(iStatorIrms100), max(iStatorIrms90), max(iStatorIrms110)))
        # document.add_paragraph('转子电流最大值为：      \t%8.3f [A],   %8.3f [A],   %8.3f [A]' % (
        #     max(iGenIrmsSJ100), max(iGenIrmsSJ90), max(iGenIrmsSJ110)))
        # document.add_paragraph(
        #     '变流器网侧电流最大值为： \t%8.3f [A],   %8.3f [A],   %8.3f [A]' % (
        #         max(iNetIrms100), max(iNetIrms90), max(iNetIrms110)))
        # document.add_paragraph(
        #     '总上网电流最大值为：    \t%8.3f [A],   %8.3f [A],   %8.3f [A]' % (
        #         max(iGridrms100), max(iGridrms90), max(iGridrms110)))
        # document.add_page_break()
        #
        # ########输出转子电压
        # document.add_heading('输出-转子电压', level=1)
        #
        # paragraph = document.add_paragraph()
        # paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        # run = paragraph.add_run("")
        # run.add_picture(r'.\fig_save\fig_voltageR.png', width=Inches(5.25))
        #
        # run = document.add_paragraph('图4. 转子线电压')
        # run.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        #
        # document.add_paragraph(r'以下三个数据分别表示工况（正常电网电压）（电网90%低电压）（电网110%高电压）')
        # document.add_paragraph(
        #     '转子线电压最大值为： \t%8.3f [V],   %8.3f [V],   %8.3f [V]' % (max(iGenVrms100), max(iGenVrms90), max(iGenVrms110)))
        # document.add_page_break()
        #
        # ######## 能力边界
        # document.add_heading('有功、无功能力边界', level=1)
        # paragraph = document.add_paragraph('定子有无功能力范围：')
        # run = paragraph.add_run("")
        # run = paragraph.add_run("根据用户界面输入的“（选填）机侧变流器最大电流”与“（选填）发电机定子最大电流”等参数，计算发电机定子侧的无功功率约束关系。")
        #
        # paragraph = document.add_paragraph()
        # paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        # run = paragraph.add_run("")
        # run.add_picture(r'.\fig_save\fig_VarAbi.png', width=Inches(5.25))
        #
        # run = document.add_paragraph('图5. 定子有无功能力边界')
        # run.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        #
        # document.add_paragraph('根据指定的定子无功功率值，计算有功输出最大值。')
        #
        # tempStr = "所有条件下，定子有功功率最大值 Pmax = "
        # tempStr = tempStr + "%.3f" % (Pmax / 1000)
        # tempStr = tempStr + "[kW]  \n"
        # tempStr = tempStr + "无功功率确定条件下 Ppoint_max = "
        # tempStr = tempStr + "%.3f" % (Puplim / 1000)
        # tempStr = tempStr + "[kW]，"
        # tempStr = tempStr + "当定子无功功率确定为 Qpoint = %.0f [kVar]" % (Qpoint / 1000)
        # document.add_paragraph(tempStr)
        #
        # document.add_paragraph('根据指定的定子有功功率值，计算无功输出最大值。')
        # tempStr = "所有条件下，定子无功功率最大值 Qmax = "
        # tempStr = tempStr + "%.3f" % (Qmax / 1000)
        # tempStr = tempStr + "[kVar]\n"
        # tempStr = tempStr + "有功功率确定条件下 Qpoint_max = "
        # tempStr = tempStr + "%.3f" % (Quplim / 1000)
        # tempStr = tempStr + "[kVar]，"
        # tempStr = tempStr + "当定子有功功率确定为 Ppoint = %.0f [kW]" % (Ppoint / 1000)
        # document.add_paragraph(tempStr)
        #
        # # document.add_paragraph()
        # document.add_page_break()
        #
        # #
        # # ax.annotate(tempStr, xy=(Qpoint, Puplim), xytext=(Qpoint-3E6, Puplim + 0), color='g',
        # #             arrowprops=dict(facecolor='green', arrowstyle="->", connectionstyle="arc3,rad=.2"),
        # #             )
        # #
        # # tempStr = "Qmax = "
        # # tempStr = tempStr + "%.3f" % (Quplim / 1000)
        # # tempStr = tempStr + "[kVar]\n"
        # # tempStr = tempStr+ "When P = %.0f [kW]"% (Ppoint / 1000)
        # # ax.annotate(tempStr, xy=(Quplim, Ppoint), xytext=(Quplim + 0.2E6, 0), color='g',
        #
        # ########结论
        # document.add_heading('结论', level=1)
        # paragraph = document.add_paragraph('根据用户在界面上输入的选填内容，进行设计合理性判断。')
        #
        #
        # table_para = document.add_table(rows=5, cols=3, style='Table Grid')
        #
        # hdr_cells = table_para.columns[0].cells
        #
        # hdr_cells[0].text = '判定项'
        # hdr_cells[1].text = '发电机转子电压'
        # hdr_cells[2].text = '网侧变流器电流'
        # hdr_cells[3].text = '机侧变流器电流'
        # hdr_cells[4].text = '发电机定子电流'
        #
        # para_cells1 = table_para.columns[1].cells
        # para_cells1[0].text = '判定标准'
        #
        # TF = [True for i in range(4)]
        # para_cells1[1].text = str(format(max(max(iGenVrms90), max(iGenVrms100), max(iGenVrms110)),
        #                                  '.2f')) + ' < ' + self.m_Vdc.GetValue() + '/1.414'
        # TF[0] = max(max(iGenVrms90), max(iGenVrms100), max(iGenVrms110)) < float(self.m_Vdc.GetValue()) / 1.414
        #
        # para_cells1[2].text = str(format(max(max(iNetIrms90), max(iNetIrms100), max(iNetIrms110)),
        #                                  '.2f')) + ' < ' + self.m_SetNetCurMax.GetValue()
        # TF[1] = max(max(iNetIrms90), max(iNetIrms100), max(iNetIrms110)) < float(self.m_SetNetCurMax.GetValue())
        #
        # para_cells1[3].text = str(format(max(max(iGenIrmsSJ90), max(iGenIrmsSJ100), max(iGenIrmsSJ110)),
        #                                  '.2f')) + ' < ' + self.m_SetGenCurMax1.GetValue()
        # TF[2] = max(max(iGenIrmsSJ90), max(iGenIrmsSJ100), max(iGenIrmsSJ110)) < float(self.m_SetGenCurMax1.GetValue())
        #
        # para_cells1[4].text = str(format(max(max(iStatorIrms90), max(iStatorIrms100), max(iStatorIrms110)),
        #                                  '.2f')) + ' < ' + self.m_SetSatMaxI.GetValue()
        # TF[3] = max(max(iStatorIrms90), max(iStatorIrms100), max(iStatorIrms110)) < float(self.m_SetSatMaxI.GetValue())
        #
        # para_cells2 = table_para.columns[2].cells
        # para_cells2[0].text = '是否通过'
        #
        # # 1,2,3,4行
        # for i in range(4):
        #     p = para_cells2[i + 1].paragraphs.pop()
        #     # print(dir(para_cells2[i + 1].paragraphs))
        #     # print(dir(p.style))
        #     # print(len(p))
        #     # p.style.font.color.rgb = RGBColor(250, 0, 0)
        #     run = p.add_run(str(TF[i]))
        #     if TF[i] is False:
        #         run.font.color.rgb = RGBColor(250, 0, 0)
        #
        # document.add_paragraph(r'')
        #
        # if all(i == True for i in TF):
        #     document.add_paragraph(r'结论:设计参数校验通过。')
        # else:
        #     paragraph = document.add_paragraph(r'结论:设计参数 %d项校验不通过！' % TF.count(False))
        #     # paragraph.style.font.color.rgb = RGBColor(250,0,0)
        #
        document.add_page_break()
        document.save('鼠笼风电机组电驱系统设计计算报告.docx')

        # ##############################################test


'''end of file'''
